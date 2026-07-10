#!/usr/bin/env python3
"""
Baut aus manifest.json (siehe render_slides.js) eine editierbare .pptx:
- Hintergrundbild (bg-N.png) füllt jede Folie randlos
- Jedes data-pptx-Textelement wird als echte, editierbare PowerPoint-Textbox
  an derselben Position/Größe/Formatierung wieder eingefügt

Usage: python build_pptx.py <manifest.json> <output.pptx>
"""
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "right": PP_ALIGN.RIGHT,
    "center": PP_ALIGN.CENTER,
    "justify": PP_ALIGN.JUSTIFY,
    "start": PP_ALIGN.LEFT,
    "end": PP_ALIGN.RIGHT,
}


def parse_css_color(css_color: str) -> RGBColor:
    """Wandelt 'rgb(r, g, b)' oder 'rgba(r, g, b, a)' aus getComputedStyle in RGBColor um."""
    m = re.match(r"rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)", css_color)
    if not m:
        return RGBColor(0, 0, 0)
    r, g, b = (int(m.group(i)) for i in (1, 2, 3))
    return RGBColor(r, g, b)


def add_background(slide, image_path: Path, slide_w: Emu, slide_h: Emu):
    slide.shapes.add_picture(str(image_path), 0, 0, width=slide_w, height=slide_h)


def add_text_element(slide, el: dict):
    left, top = Inches(el["x"]), Inches(el["y"])
    width, height = Inches(max(el["w"], 0.05)), Inches(max(el["h"], 0.05))

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    # Kein Auto-Shrink: wir vertrauen der HTML-Quelle als Wahrheit für die Größe.
    from pptx.enum.text import MSO_AUTO_SIZE
    tf.auto_size = MSO_AUTO_SIZE.NONE

    paragraphs = el["paragraphs"]
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ALIGN_MAP.get(para.get("align", "left"), PP_ALIGN.LEFT)
        if para.get("lineHeightPt"):
            # line_spacing als Pt-Wert entspricht der Zeilenhöhe direkt (python-pptx erlaubt Pt oder Faktor)
            try:
                p.line_spacing = Pt(para["lineHeightPt"])
            except Exception:
                pass
        if para.get("spaceAfterPt"):
            p.space_after = Pt(para["spaceAfterPt"])

        runs = list(para["runs"])
        if para.get("isListItem") and runs:
            bullet = "• " if para.get("listType") == "ul" else f"{i + 1}. "
            runs = [dict(runs[0], text=bullet + runs[0]["text"])] + runs[1:]

        for run_data in runs:
            r = p.add_run()
            text = run_data["text"]
            style = run_data["style"]
            r.text = text.upper() if style.get("uppercase") else text
            r.font.size = Pt(style["fontSizePt"])
            r.font.bold = bool(style.get("bold"))
            r.font.italic = bool(style.get("italic"))
            r.font.underline = bool(style.get("underline"))
            if style.get("fontFamily"):
                r.font.name = style["fontFamily"]
            r.font.color.rgb = parse_css_color(style.get("color", "rgb(0,0,0)"))


def main():
    if len(sys.argv) != 3:
        print("Usage: python build_pptx.py <manifest.json> <output.pptx>")
        sys.exit(1)

    manifest_path = Path(sys.argv[1])
    output_path = Path(sys.argv[2])
    base_dir = manifest_path.parent

    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))

    prs = Presentation()
    prs.slide_width = Inches(manifest["slideWidthIn"])
    prs.slide_height = Inches(manifest["slideHeightIn"])
    blank_layout = prs.slide_layouts[6]

    for slide_data in manifest["slides"]:
        slide = prs.slides.add_slide(blank_layout)
        add_background(slide, base_dir / slide_data["bg"], prs.slide_width, prs.slide_height)
        for el in slide_data["textElements"]:
            add_text_element(slide, el)
        print(f"Folie {slide_data['slideNum']}: {len(slide_data['textElements'])} Textboxen eingefügt")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Gespeichert: {output_path}")


if __name__ == "__main__":
    main()
